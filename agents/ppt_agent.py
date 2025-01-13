from typing import Dict, Any
from .base_agent import BaseAgent
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

class PPTAgent(BaseAgent):
    def __init__(self):
        super().__init__()  # PPT agent doesn't need LLM
        
    def process(self, state: Dict[str, Any]) -> Dict[str, Any]:
        """Create a presentation from the analysis results"""
        try:
            response = state['response']
            query = state['query']
            
            # Create presentation
            prs = Presentation()
            
            # Set slide size to 16:9
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            
            # Add title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = "Analysis Report"
            title_slide.placeholders[1].text = query
            
            # Split content into sections
            sections = self._split_content(response)
            
            # Add content slides
            for section in sections:
                content_slide = prs.slides.add_slide(prs.slide_layouts[1])
                content_slide.shapes.title.text = section['title']
                content = content_slide.placeholders[1]
                content.text = section['content']
                
                # Format text
                for paragraph in content.text_frame.paragraphs:
                    paragraph.font.size = Pt(18)
            
            # Save presentation
            output_path = "output/presentation.pptx"
            prs.save(output_path)
            state['presentation_path'] = output_path
            
            return state
            
        except Exception as e:
            print(f"Error in PPTAgent: {str(e)}")
            raise

    def _split_content(self, content: str) -> list:
        """Split content into presentation-friendly sections"""
        try:
            # Split by paragraphs
            paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
            sections = []
            
            for i, para in enumerate(paragraphs, 1):
                # Limit content length per slide
                if len(para) > 500:
                    para = para[:497] + '...'
                
                sections.append({
                    'title': f'Key Point {i}',
                    'content': para
                })
            
            return sections
        except Exception:
            # Fallback to single section if splitting fails
            return [{
                'title': 'Analysis Results',
                'content': content[:1000]  # Limit content if needed
            }]

    def validate_input(self, state: Dict[str, Any]) -> bool:
        """Validate required state inputs"""
        return all(k in state for k in ['query', 'response']) 